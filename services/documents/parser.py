"""Document parsing with defense-in-depth limits.

All parser/library failures are normalized to safe structured AppErrors;
stack traces, local paths, XML internals, and library messages never reach
the client. The processing timeout is cooperative (checked between
pages/slides) because parsing runs in-process. See docs/DEPLOYMENT.md for
the threat model and operational guidance.
"""
import io
import logging
import time
import zipfile
from pathlib import Path, PurePosixPath

from docx import Document as DocxDocument
from pptx import Presentation
from pypdf import PdfReader

from utils.errors import AppError

logger = logging.getLogger(__name__)

OFFICE_CONTENT = {"docx": "word/", "pptx": "ppt/"}


def inspect_office_archive(blob, ext, max_files, max_expanded):
    """Reject encrypted entries, traversal names, and archive bombs."""
    try:
        with zipfile.ZipFile(blob) as archive:
            infos = archive.infolist()
            if len(infos) > max_files:
                raise AppError("ARCHIVE_BOMB", "The document contains too many archive entries.", 422)
            total = 0
            for info in infos:
                name = PurePosixPath(info.filename)
                if info.flag_bits & 1 or name.is_absolute() or ".." in name.parts:
                    raise AppError("INVALID_ARCHIVE", "Encrypted or unsafe archive entries are not allowed.", 422)
                total += info.file_size
                if total > max_expanded or (info.compress_size and info.file_size / max(1, info.compress_size) > 200):
                    raise AppError("ARCHIVE_BOMB", "The document expands beyond safe limits.", 422)
            names = {x.filename for x in infos}
            if "[Content_Types].xml" not in names or not any(x.startswith(OFFICE_CONTENT[ext]) for x in names):
                raise AppError("INVALID_FILE_SIGNATURE", "Office document structure is invalid.", 415)
    except zipfile.BadZipFile as exc:
        raise AppError("INVALID_FILE_SIGNATURE", "Office document structure is invalid.", 415) from exc


def parse_document(path, *, max_pages=250, max_chars=500000, max_files=1000, max_expanded=80 * 1024 * 1024, timeout=30):
    """Parse PDF/DOCX/PPTX/TXT into (text, page_map).

    The file is read once into memory and all parsers operate on BytesIO so
    every file handle is deterministically closed.
    """
    path = Path(path)
    ext = path.suffix.lower()
    start = time.monotonic()
    pages = []

    def add(number, text):
        if len(pages) >= max_pages:
            raise AppError("DOCUMENT_PAGE_LIMIT", "The document has too many pages or slides.", 422)
        if time.monotonic() - start > timeout:
            raise AppError("DOCUMENT_TIMEOUT", "Document processing exceeded the time limit.", 422)
        pages.append({"page": number, "text": str(text)[:max_chars]})

    try:
        raw = path.read_bytes()
        if ext == ".pdf":
            reader = PdfReader(io.BytesIO(raw), strict=True)
            if reader.is_encrypted:
                raise AppError("DOCUMENT_ENCRYPTED", "Password-protected documents are not supported.", 422)
            if len(reader.pages) > max_pages:
                raise AppError("DOCUMENT_PAGE_LIMIT", "The PDF has too many pages.", 422)
            for i, page in enumerate(reader.pages, 1):
                add(i, page.extract_text() or "")
        elif ext in {".docx", ".pptx"}:
            blob = io.BytesIO(raw)
            inspect_office_archive(io.BytesIO(raw), ext[1:], max_files, max_expanded)
            if ext == ".docx":
                document = DocxDocument(blob)
                add(1, "\n".join(p.text for p in document.paragraphs))
            else:
                presentation = Presentation(blob)
                if len(presentation.slides) > max_pages:
                    raise AppError("DOCUMENT_PAGE_LIMIT", "The presentation has too many slides.", 422)
                for i, slide in enumerate(presentation.slides, 1):
                    add(i, "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")))
        elif ext == ".txt":
            if b"\x00" in raw[:4096]:
                raise AppError("INVALID_TEXT_FILE", "Binary data is not accepted as a text document.", 415)
            add(1, raw.decode("utf-8", errors="strict"))
        else:
            raise AppError("UNSUPPORTED_FILE", "Allowed types: PDF, DOCX, PPTX, TXT.", 415)
    except AppError:
        raise
    except Exception as exc:  # parser libraries raise many exception types
        logger.warning("Document parse failed (%s): %s", ext, type(exc).__name__)
        raise AppError("DOCUMENT_PARSE_FAILED", "The document could not be parsed safely.", 422) from exc

    text = "\n\n".join(f"[Document page {p['page']}]\n{p['text']}" for p in pages)[:max_chars]
    if not text.strip() or not any(p["text"].strip() for p in pages):
        raise AppError("DOCUMENT_EMPTY", "No readable text was found in the document.", 422)
    return text, [{"page": p["page"], "characters": len(p["text"])} for p in pages]
